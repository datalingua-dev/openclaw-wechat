import sys
import os
import argparse
import traceback

def extract_pdf(path):
    import pdfplumber
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text.append(t)
    return "\n\n".join(text)

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        text.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_html(path):
    from bs4 import BeautifulSoup
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'html.parser')
        # 移除脚本和样式
        for script in soup(["script", "style"]):
            script.extract()
        return soup.get_text(separator='\n', strip=True)

def extract_text(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".doc": extract_docx, # docx library partially supports some .doc via compatibility or may fail
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
}

def main():
    parser = argparse.ArgumentParser(description="Extract text from various file formats")
    parser.add_argument("file_path", help="Path to the file to extract text from")
    parser.add_argument("--max_chars", type=int, default=50000, help="Maximum number of characters to extract")
    args = parser.parse_args()

    file_path = args.file_path
    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext in EXTRACTORS:
            extractor = EXTRACTORS[ext]
            content = extractor(file_path)
        else:
            # Default to plain text
            content = extract_text(file_path)
        
        if content:
            if len(content) > args.max_chars:
                content = content[:args.max_chars] + "\n\n...(内容已截断)"
            print(content)
        else:
            print("", end="")
    except Exception as e:
        print(f"Error extracting {file_path}: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
